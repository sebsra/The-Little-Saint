"""
Project Structure Scanner

A tool that scans directories and creates a detailed report of project structure,
including directory hierarchy and file contents. Useful for documenting project
organization or sharing code structure.

Scans text files with configurable inclusion or exclusion lists and supports document formats.
"""

import os
import sys
from pathlib import Path
from typing import List, Set, Optional, TextIO, Dict, Any
from datetime import datetime
from dataclasses import dataclass, field
from enum import Enum

# ====================================================================
# DEFAULT CONFIGURATION CONSTANTS
# ====================================================================

# Document processing defaults
DEFAULT_DOC_SUPPORT_ENABLED = True

# File filtering mode
class FilterMode(Enum):
    EXCLUDE = "exclude"  # Process all files except those with excluded extensions
    INCLUDE = "include"  # Process only files with included extensions

DEFAULT_FILTER_MODE = FilterMode.EXCLUDE

# File exclusion defaults
DEFAULT_EXCLUDED_EXTENSIONS = {
    '.exe', '.dll', '.so', '.pyc', '.pyd', '.obj', '.jpg', '.jpeg', '.png', '.gif', 
    '.bmp', '.ico', '.zip', '.tar', '.gz', '.7z', '.rar', '.bin', '.dat', '.db', 
    '.sqlite', '.mp3', '.mp4', '.wav'
}


# File inclusion defaults (used when in INCLUDE mode)
DEFAULT_INCLUDED_EXTENSIONS = {".gd", ".godot",}

# Default specific file inclusion/exclusion
DEFAULT_EXCLUDED_FILES = set()
DEFAULT_INCLUDED_FILES = set()  # No specific files included by default

# Directory exclusion defaults
DEFAULT_EXCLUDED_DIRS = {
    "__pycache__", ".venv", "venv", ".git", "node_modules", "build", "dist", 
    ".idea", ".vs", ".vscode", ".pytest_cache", 
}

# Content display defaults
DEFAULT_MAX_FILE_SIZE_KB = 1024  # 1MB limit
DEFAULT_MAX_LINE_COUNT = 1000    # 1000 lines per file


@dataclass
class DocumentSupport:
    """Document type support configuration."""
    # Document type settings
    pdf_enabled: bool
    word_enabled: bool
    excel_enabled: bool
    powerpoint_enabled: bool
    
    # Availability flags (determined at runtime)
    pdf_available: bool = field(default=False, init=False)
    word_available: bool = field(default=False, init=False)
    excel_available: bool = field(default=False, init=False)
    powerpoint_available: bool = field(default=False, init=False)
    
    def __post_init__(self):
        """Check which document libraries are available."""
        # Check PDF support
        try:
            from PyPDF2 import PdfReader
            self.pdf_available = True
        except ImportError:
            if self.pdf_enabled:
                print("Warning: PDF support enabled but PyPDF2 library not found.")
            
        # Check Word support
        try:
            import docx
            self.word_available = True
        except ImportError:
            if self.word_enabled:
                print("Warning: Word support enabled but python-docx library not found.")
            
        # Check Excel support
        try:
            import openpyxl
            self.excel_available = True
        except ImportError:
            if self.excel_enabled:
                print("Warning: Excel support enabled but openpyxl library not found.")
            
        # Check PowerPoint support
        try:
            from pptx import Presentation
            self.powerpoint_available = True
        except ImportError:
            if self.powerpoint_enabled:
                print("Warning: PowerPoint support enabled but python-pptx library not found.")
    
    def get_status_report(self) -> str:
        """Generate a status report of document support."""
        status_lines = ["Document Support Status:"]
        
        status_lines.append(f"- PDF: {'Enabled' if self.pdf_enabled else 'Disabled'} "
                          f"({'Available' if self.pdf_available else 'Not available - install PyPDF2'})")
        
        status_lines.append(f"- Word: {'Enabled' if self.word_enabled else 'Disabled'} "
                          f"({'Available' if self.word_available else 'Not available - install python-docx'})")
        
        status_lines.append(f"- Excel: {'Enabled' if self.excel_enabled else 'Disabled'} "
                          f"({'Available' if self.excel_available else 'Not available - install openpyxl'})")
        
        status_lines.append(f"- PowerPoint: {'Enabled' if self.powerpoint_enabled else 'Disabled'} "
                          f"({'Available' if self.powerpoint_available else 'Not available - install python-pptx'})")
        
        return "\n".join(status_lines)


class ScannerConfigError(Exception):
    """Exception raised for configuration errors in the scanner."""
    pass


@dataclass
class ScannerConfig:
    """Configuration for the project scanner."""
    # Required parameters
    directory: Path
    output_file_path: Path
    
    # File filtering options
    filter_mode: FilterMode
    excluded_extensions: Set[str]
    included_extensions: Set[str]
    excluded_files: Set[str]  # Specific files to exclude by name
    included_files: Set[str]  # Specific files to include by name (overrides extension rules)
    
    # Path blacklisting
    blacklisted_paths: Set[Path]
    
    # Directory exclusion settings
    excluded_dirs: Set[str]
    
    # Content display settings
    show_structure: bool
    show_contents: bool
    
    # File content limits
    max_file_size_kb: Optional[int]
    max_line_count: Optional[int]
    
    # Document support
    doc_support: DocumentSupport
    
    def __post_init__(self):
        """Validate and normalize the configuration."""
        # Validate directory
        if not self.directory.exists():
            raise ScannerConfigError(f"Directory does not exist: {self.directory}")
        if not self.directory.is_dir():
            raise ScannerConfigError(f"Not a directory: {self.directory}")
            
        # Ensure output directory exists
        output_dir = self.output_file_path.parent
        if not output_dir.exists():
            raise ScannerConfigError(f"Output directory does not exist: {output_dir}")
            
        # Validate content limits
        if self.max_file_size_kb is not None and self.max_file_size_kb <= 0:
            raise ScannerConfigError("max_file_size_kb must be positive if specified")
            
        if self.max_line_count is not None and self.max_line_count <= 0:
            raise ScannerConfigError("max_line_count must be positive if specified")
            
        # Normalize blacklisted paths to absolute paths
        normalized_blacklist = set()
        for path in self.blacklisted_paths:
            if isinstance(path, str):
                path = Path(path)
            normalized_blacklist.add(path.absolute())
        self.blacklisted_paths = normalized_blacklist
        
        # Validate filter mode
        if not isinstance(self.filter_mode, FilterMode):
            raise ScannerConfigError(f"Invalid filter mode: {self.filter_mode}")
            
        # If in include mode, we should have included extensions
        if self.filter_mode == FilterMode.INCLUDE and not self.included_extensions:
            raise ScannerConfigError("When using INCLUDE filter mode, included_extensions cannot be empty")


class ScannerError(Exception):
    """Exception raised for errors during scanning operations."""
    pass


class Scanner:
    """Scans project structure and generates reports."""
    
    def __init__(self, config: ScannerConfig):
        """Initialize the scanner with configuration."""
        self.config = config
    
    def scan(self) -> None:
        """Perform the full project scan."""
        try:
            with open(self.config.output_file_path, 'w', encoding='utf-8') as output_file:
                self._write_header(output_file)
                
                # Collect files if showing structure
                file_list = []
                if self.config.show_structure:
                    output_file.write("\nDIRECTORY STRUCTURE\n")
                    output_file.write("-" * 18 + "\n\n")
                    file_list = self._scan_directory(self.config.directory, output_file=output_file)
                
                # Process file contents if enabled
                if self.config.show_contents:
                    output_file.write("\nFILE CONTENTS\n")
                    output_file.write("-" * 13 + "\n\n")
                    self._process_files(file_list, output_file)
                else:
                    output_file.write("\nFile contents skipped\n")
                    
                output_file.write(f"\n{'=' * 50}\n")
                output_file.write(f"Scan completed at {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
        except PermissionError as e:
            raise ScannerError(f"Permission denied when writing to output file: {e}")
        except IOError as e:
            raise ScannerError(f"I/O error when writing to output file: {e}")
    
    def _write_header(self, output_file: TextIO) -> None:
        """Write the scan report header."""
        output_file.write(f"PROJECT STRUCTURE SCAN\n")
        output_file.write(f"{'=' * 50}\n")
        output_file.write(f"Directory: {self.config.directory.absolute()}\n")
        output_file.write(f"Date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
        
        # Show filter mode and relevant extension list
        output_file.write(f"Filter mode: {self.config.filter_mode.value}\n")
        if self.config.filter_mode == FilterMode.EXCLUDE:
            output_file.write(f"Excluded file types: {', '.join(sorted(self.config.excluded_extensions))}\n")
        else:  # INCLUDE mode
            output_file.write(f"Included file types: {', '.join(sorted(self.config.included_extensions))}\n")
        
        # Show excluded and included files
        if self.config.excluded_files:
            output_file.write(f"Excluded files: {', '.join(sorted(self.config.excluded_files))}\n")
            
        if self.config.included_files:
            output_file.write(f"Included files: {', '.join(sorted(self.config.included_files))}\n")
            
        if self.config.blacklisted_paths:
            output_file.write(f"Blacklisted paths: {', '.join(str(p) for p in self.config.blacklisted_paths)}\n")
            
        output_file.write(f"\n{self.config.doc_support.get_status_report()}\n")
        output_file.write(f"{'=' * 50}")
    
    def _scan_directory(self, directory: Path, indent: str = "", output_file: TextIO = None, 
                       file_list: Optional[List[Path]] = None) -> List[Path]:
        """
        Recursively scan a directory and collect files.
        
        Args:
            directory: The directory to scan
            indent: Current indentation level for output formatting
            output_file: File to write output to
            file_list: List to collect file paths
            
        Returns:
            List of file paths found during scanning
        """
        if file_list is None:
            file_list = []
            
        # Skip if this directory is blacklisted
        if directory.absolute() in self.config.blacklisted_paths:
            output_file.write(f"{indent}[DIR] {directory.name} (blacklisted)\n")
            return file_list
            
        try:
            items = sorted(os.listdir(directory))
            
            for item in items:
                item_path = directory / item
                
                # Skip if this path is blacklisted
                if item_path.absolute() in self.config.blacklisted_paths:
                    output_file.write(f"{indent}{'- ' if item_path.is_file() else '[DIR] '}{item} (blacklisted)\n")
                    continue
                
                if item_path.is_dir():
                    # Skip excluded directories
                    if item in self.config.excluded_dirs or item.startswith('.'):
                        output_file.write(f"{indent}[DIR] {item} (excluded)\n")
                        continue
                        
                    output_file.write(f"{indent}[DIR] {item}\n")
                    self._scan_directory(item_path, indent + "  ", output_file, file_list)
                else:
                    # Check file extension based on filter mode
                    extension = item_path.suffix.lower()
                    skip_file = False
                    
                    if self.config.filter_mode == FilterMode.EXCLUDE:
                        # Skip files with excluded extensions
                        if extension in self.config.excluded_extensions:
                            output_file.write(f"{indent}- {item} (excluded by extension)\n")
                            skip_file = True
                    else:  # INCLUDE mode
                        # Skip files without included extensions
                        if extension not in self.config.included_extensions:
                            #output_file.write(f"{indent}- {item} (not in included extensions)\n")
                            skip_file = True
                    
                    if not skip_file:
                        output_file.write(f"{indent}- {item}\n")
                        file_list.append(item_path)
                    
        except PermissionError:
            output_file.write(f"{indent}[Permission Denied]\n")
        except OSError as e:
            output_file.write(f"{indent}[Error: {str(e)}]\n")
            
        return file_list
    
    def _should_process_file(self, file_path: Path) -> bool:
        """
        Determine if a file should be processed based on filter mode, extensions, and excluded/included files.
        
        Args:
            file_path: Path to the file to check
            
        Returns:
            True if the file should be processed, False otherwise
        """
        # First check if the file is specifically included by name (overrides other rules)
        if file_path.name in self.config.included_files:
            return True
            
        # Then check if the file is specifically excluded by name
        if file_path.name in self.config.excluded_files:
            return False
            
        extension = file_path.suffix.lower()
        
        if self.config.filter_mode == FilterMode.EXCLUDE:
            # Process all files except those with excluded extensions
            return extension not in self.config.excluded_extensions
        else:  # INCLUDE mode
            # Process only files with included extensions
            return extension in self.config.included_extensions
    
    def _process_files(self, file_list: List[Path], output_file: TextIO) -> None:
        """
        Process and display file contents.
        
        Args:
            file_list: List of files to process
            output_file: File to write output to
        """
        # Filter out special files and the output file itself
        script_path = Path(__file__).absolute()
        output_path = self.config.output_file_path.absolute()
        
        filtered_files = [
            f for f in file_list 
            if f.absolute() != output_path and f.absolute() != script_path
        ]
        
        # Process each file
        for file_path in filtered_files:
            try:
                # Apply filter based on mode and extensions
                if not self._should_process_file(file_path):
                    continue
                
                # Check file size if max size is specified
                if self.config.max_file_size_kb is not None:
                    file_size_kb = file_path.stat().st_size / 1024
                    if file_size_kb > self.config.max_file_size_kb:
                        output_file.write(
                            f"\n[Skipped {file_path}: Size {file_size_kb:.1f}KB exceeds limit of "
                            f"{self.config.max_file_size_kb}KB]\n"
                        )
                        continue
                
                # Process the file based on its type
                if self._is_document_file(file_path):
                    self._process_document(file_path, output_file)
                elif not self._is_binary_file(file_path):
                    self._process_text_file(file_path, output_file)
                    
            except FileNotFoundError:
                output_file.write(f"\n[Error: File not found: {file_path}]\n")
            except PermissionError:
                output_file.write(f"\n[Error: Permission denied: {file_path}]\n")
            except UnicodeDecodeError:
                output_file.write(f"\n[Error: Unable to decode file {file_path} - not a valid text file]\n")
            except Exception as e:
                output_file.write(f"\n[Error processing file {file_path}: {e}]\n")
    
    def _is_document_file(self, file_path: Path) -> bool:
        """Check if a file is a supported document type."""
        suffix = file_path.suffix.lower()
        doc = self.config.doc_support
        
        return (
            (suffix == '.pdf' and doc.pdf_enabled and doc.pdf_available) or
            (suffix in ('.docx', '.doc') and doc.word_enabled and doc.word_available) or
            (suffix in ('.xlsx', '.xls') and doc.excel_enabled and doc.excel_available) or
            (suffix in ('.pptx', '.ppt') and doc.powerpoint_enabled and doc.powerpoint_available)
        )
    
    def _is_binary_file(self, file_path: Path) -> bool:
        """Check if a file appears to be binary."""
        # If we're in exclude mode, check if it's a known binary extension
        if self.config.filter_mode == FilterMode.EXCLUDE:
            if file_path.suffix.lower() in self.config.excluded_extensions:
                return True
            
        # Check for binary indicators
        try:
            with open(file_path, 'rb') as file:
                chunk = file.read(4096)
                
                # Check for NULL bytes (common in binary files)
                if b'\x00' in chunk:
                    return True
                    
                # Calculate percentage of printable ASCII characters
                printable = sum(32 <= b <= 126 or b in (9, 10, 13) for b in chunk)
                if chunk and printable / len(chunk) < 0.7:
                    return True
                    
                return False
        except Exception as e:
            # If we can't read it, report the error and assume binary to be safe
            print(f"Warning: Error reading file {file_path}: {e}")
            return True
    
    def _process_text_file(self, file_path: Path, output_file: TextIO) -> None:
        """Process and display text file contents."""
        output_file.write(f"\n{'=' * 40}\n")
        output_file.write(f"Contents of {file_path}:\n")
        output_file.write(f"{'=' * 40}\n")
        
        try:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
                # Read lines with limit if specified
                if self.config.max_line_count is not None:
                    lines = []
                    for i, line in enumerate(file):
                        if i >= self.config.max_line_count:
                            output_file.write("".join(f"  {l}" for l in lines))
                            output_file.write(f"\n... (truncated, {self.config.max_line_count} of {i+1}+ lines shown) ...\n")
                            break
                        lines.append(line)
                    else:
                        # This executes if the loop completed without breaking
                        output_file.write("".join(f"  {l}" for l in lines))
                else:
                    # No line limit
                    content = file.read()
                    output_file.write("".join(f"  {l}" for l in content.splitlines(True)))
        except UnicodeDecodeError:
            raise ScannerError(f"File contains non-UTF-8 characters or is binary: {file_path}")
    
    def _process_document(self, file_path: Path, output_file: TextIO) -> None:
        """Process document files (PDF, Word, Excel, PowerPoint)."""
        doc_type = file_path.suffix.upper()[1:]  # Get extension without dot, uppercase
        
        output_file.write(f"\n{'=' * 40}\n")
        output_file.write(f"Contents of {file_path} ({doc_type}):\n")
        output_file.write(f"{'=' * 40}\n")
        
        try:
            doc_text = self._extract_document_text(file_path)
            
            # Apply line limits if specified
            if self.config.max_line_count is not None:
                lines = doc_text.splitlines()
                if len(lines) > self.config.max_line_count:
                    limited_lines = lines[:self.config.max_line_count]
                    output_file.write("\n".join(f"  {line}" for line in limited_lines))
                    output_file.write(f"\n... (truncated, {self.config.max_line_count} of {len(lines)} lines shown) ...\n")
                else:
                    output_file.write("\n".join(f"  {line}" for line in lines))
            else:
                output_file.write("\n".join(f"  {line}" for line in doc_text.splitlines()))
        except Exception as e:
            raise ScannerError(f"Error extracting document content from {file_path}: {e}")
    
    def _extract_document_text(self, file_path: Path) -> str:
        """Extract text from document files based on their type."""
        suffix = file_path.suffix.lower()
        doc = self.config.doc_support
        
        # Handle PDF files
        if suffix == '.pdf' and doc.pdf_enabled and doc.pdf_available:
            return self._extract_pdf_text(file_path)
            
        # Handle Word documents
        if suffix in ('.docx', '.doc') and doc.word_enabled and doc.word_available:
            return self._extract_word_text(file_path)
            
        # Handle Excel files
        if suffix in ('.xlsx', '.xls') and doc.excel_enabled and doc.excel_available:
            return self._extract_excel_text(file_path)
            
        # Handle PowerPoint files
        if suffix in ('.pptx', '.ppt') and doc.powerpoint_enabled and doc.powerpoint_available:
            return self._extract_powerpoint_text(file_path)
            
        raise ScannerError(f"Unsupported document type: {suffix}")
    
    def _extract_pdf_text(self, file_path: Path) -> str:
        """Extract text from a PDF file."""
        try:
            from PyPDF2 import PdfReader
            
            text_parts = []
            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                num_pages = len(pdf_reader.pages)
                
                for page_num in range(num_pages):
                    page = pdf_reader.pages[page_num]
                    text_parts.append(f"--- Page {page_num + 1} of {num_pages} ---")
                    page_text = page.extract_text() or "[No text content on this page]"
                    text_parts.append(page_text)
                    
            return '\n'.join(text_parts)
        except Exception as e:
            raise ScannerError(f"Error extracting text from PDF: {str(e)}")
    
    def _extract_word_text(self, file_path: Path) -> str:
        """Extract text from a Word document."""
        try:
            import docx
            
            doc = docx.Document(file_path)
            text_parts = []
            
            # Extract document properties
            text_parts.append(f"--- Document Properties ---")
            props = doc.core_properties
            text_parts.append(f"Title: {props.title or 'Not set'}")
            text_parts.append(f"Author: {props.author or 'Not set'}")
            text_parts.append(f"Created: {props.created or 'Unknown'}")
            text_parts.append(f"Modified: {props.modified or 'Unknown'}")
            
            # Extract document content
            text_parts.append(f"\n--- Document Content ---")
            for i, para in enumerate(doc.paragraphs):
                if para.text.strip():  # Skip empty paragraphs
                    text_parts.append(f"Paragraph {i+1}: {para.text}")
            
            # Extract tables
            tables_count = len(doc.tables)
            if tables_count > 0:
                text_parts.append(f"\n--- Tables ({tables_count}) ---")
                for i, table in enumerate(doc.tables):
                    text_parts.append(f"Table {i+1}:")
                    for row in table.rows:
                        row_text = ' | '.join(cell.text for cell in row.cells)
                        text_parts.append(f"  {row_text}")
            
            return '\n'.join(text_parts)
        except Exception as e:
            raise ScannerError(f"Error extracting text from Word document: {str(e)}")
    
    def _extract_excel_text(self, file_path: Path) -> str:
        """Extract text from an Excel file."""
        try:
            import openpyxl
            
            workbook = openpyxl.load_workbook(file_path, data_only=True)
            text_parts = []
            
            text_parts.append(f"--- Workbook Properties ---")
            text_parts.append(f"Sheets: {', '.join(workbook.sheetnames)}")
            
            # Process each worksheet
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text_parts.append(f"\n--- Sheet: {sheet_name} ---")
                
                # Get sheet dimensions
                min_row, min_col = 1, 1
                max_row = max(1, sheet.max_row)
                max_col = max(1, sheet.max_column)
                
                # Extract headers (first row)
                headers = []
                for col in range(min_col, max_col + 1):
                    cell = sheet.cell(min_row, col)
                    headers.append(str(cell.value) if cell.value is not None else "")
                
                text_parts.append("Headers: " + " | ".join(headers))
                text_parts.append("Data:")
                
                # Extract data rows (limited to reasonable amount)
                max_rows_to_show = min(100, max_row)
                for row in range(min_row + 1, max_rows_to_show + 1):
                    row_data = []
                    for col in range(min_col, max_col + 1):
                        cell = sheet.cell(row, col)
                        row_data.append(str(cell.value) if cell.value is not None else "")
                    text_parts.append("  " + " | ".join(row_data))
                
                if max_row > max_rows_to_show:
                    text_parts.append(f"  ... (showing {max_rows_to_show} of {max_row} rows) ...")
            
            return '\n'.join(text_parts)
        except Exception as e:
            raise ScannerError(f"Error extracting text from Excel file: {str(e)}")
    
    def _extract_powerpoint_text(self, file_path: Path) -> str:
        """Extract text from a PowerPoint file."""
        try:
            from pptx import Presentation
            
            presentation = Presentation(file_path)
            text_parts = []
            
            text_parts.append(f"--- Presentation Properties ---")
            text_parts.append(f"Slides: {len(presentation.slides)}")
            
            # Process each slide
            for i, slide in enumerate(presentation.slides):
                text_parts.append(f"\n--- Slide {i+1} ---")
                
                # Get slide title
                if slide.shapes.title:
                    text_parts.append(f"Title: {slide.shapes.title.text}")
                
                # Extract text from all shapes
                shape_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        shape_texts.append(shape.text)
                
                if shape_texts:
                    text_parts.append("Content:")
                    for text in shape_texts:
                        text_parts.append(f"  {text}")
                else:
                    text_parts.append("  [No text content on this slide]")
            
            return '\n'.join(text_parts)
        except Exception as e:
            raise ScannerError(f"Error extracting text from PowerPoint file: {str(e)}")


def create_default_config(directory_path: str, output_file_path: str, 
                         filter_mode: FilterMode = DEFAULT_FILTER_MODE) -> ScannerConfig:
    """Create a scanner configuration with default values."""
    
    # Create document support with explicit settings
    doc_support = DocumentSupport(
        pdf_enabled=DEFAULT_DOC_SUPPORT_ENABLED,
        word_enabled=DEFAULT_DOC_SUPPORT_ENABLED,
        excel_enabled=DEFAULT_DOC_SUPPORT_ENABLED,
        powerpoint_enabled=DEFAULT_DOC_SUPPORT_ENABLED
    )
    
    # Create and validate scanner configuration
    config = ScannerConfig(
        directory=Path(directory_path),
        output_file_path=Path(output_file_path),
        filter_mode=filter_mode,
        excluded_extensions=DEFAULT_EXCLUDED_EXTENSIONS,
        included_extensions=DEFAULT_INCLUDED_EXTENSIONS,
        excluded_files=DEFAULT_EXCLUDED_FILES.copy(),
        included_files=DEFAULT_INCLUDED_FILES.copy(),
        blacklisted_paths=set(),
        excluded_dirs=DEFAULT_EXCLUDED_DIRS,
        show_structure=True,
        show_contents=True,
        max_file_size_kb=DEFAULT_MAX_FILE_SIZE_KB,
        max_line_count=DEFAULT_MAX_LINE_COUNT,
        doc_support=doc_support
    )
    
    return config

def scan_core_scripts(base_dir: str, output_file_path: str) -> None:
    """
    Scan the core script directories in inclusion mode and combine the results.
    
    Args:
        base_dir: Base directory path
        output_file_path: Path for the final combined output file
    """
    import tempfile
    import shutil
    
    # Core paths to scan
    core_paths = [
        "scripts/core/projectiles",
        "scripts/core/base_classes",
        "scripts/core/enemies"
    ]
    
    # Create Path objects with proper OS-specific path handling
    base_path = Path(base_dir)
    temp_files = []
    
    try:
        print(f"Scanning core script directories...")
        
        # Scan each directory and save to a temporary file
        for i, rel_path in enumerate(core_paths):
            path = base_path / rel_path
            
            # Skip if directory doesn't exist
            if not path.exists():
                print(f"Warning: Path doesn't exist: {path}")
                continue
                
            # Create temporary file for this scan
            with tempfile.NamedTemporaryFile(delete=False, suffix='.txt') as tmp:
                temp_file_path = tmp.name
                temp_files.append(temp_file_path)
            
            # Create scanner config in INCLUDE mode
            config = create_default_config(
                str(path), 
                temp_file_path,
                filter_mode=FilterMode.INCLUDE
            )
            
            # Set included extensions to focus on script files
            config.included_extensions = {".gd", ".cs", ".gdshader", ".shader", ".json", ".cfg", ".tscn", ".tres"}
            
            # Run the scanner
            print(f"Scanning {rel_path}...")
            scanner = Scanner(config)
            scanner.scan()
        
        # Combine all temporary files into the final output
        with open(output_file_path, 'w', encoding='utf-8') as final_output:
            # Write header
            final_output.write(f"CORE SCRIPTS SCAN REPORT\n")
            final_output.write(f"{'=' * 50}\n")
            final_output.write(f"Date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            final_output.write(f"Scanned directories:\n")
            for path in core_paths:
                final_output.write(f"- {path}\n")
            final_output.write(f"{'=' * 50}\n\n")
            
            # Append content from each temp file
            for i, temp_file in enumerate(temp_files):
                if not Path(temp_file).exists():
                    continue
                    
                # Add section header for this part
                if i < len(core_paths):
                    final_output.write(f"\n{'#' * 30}\n")
                    final_output.write(f"# {core_paths[i]}\n")
                    final_output.write(f"{'#' * 30}\n\n")
                
                # Copy content (skip the header section)
                with open(temp_file, 'r', encoding='utf-8') as tmp:
                    # Skip header lines
                    header_done = False
                    for line in tmp:
                        if not header_done:
                            if line.startswith('====='):
                                header_done = True
                            continue
                        final_output.write(line)
        
        print(f"Core scripts scan completed. Combined report saved to {output_file_path}")
    
    finally:
        # Clean up temporary files
        for temp_file in temp_files:
            try:
                if os.path.exists(temp_file):
                    os.unlink(temp_file)
            except Exception as e:
                print(f"Warning: Failed to delete temporary file {temp_file}: {e}")

def main() -> int:
    """
    Main entry point for the scanner.
    Configure all settings here and run the scanner.
    """
    try:
        # ====================================================================
        # CONFIGURATION
        # ====================================================================
        
        import argparse
        
        parser = argparse.ArgumentParser(description="Project Structure Scanner")
        parser.add_argument("directory", nargs="?", default=".", help="Directory to scan (default: current directory)")
        parser.add_argument("output", nargs="?", default="project_as_text.txt", help="Output file path")
        parser.add_argument("--include", "-i", action="store_true", help="Use inclusion mode (only scan specified extensions)")
        parser.add_argument("--extensions", "-e", type=str, help="Comma-separated list of extensions to include/exclude")
        parser.add_argument("--exclude-files", "-xf", type=str, 
                          help="Comma-separated list of specific files to exclude (e.g., 'base_level.tscn,player.gd')")
        parser.add_argument("--include-files", "-if", type=str,
                          help="Comma-separated list of specific files to include regardless of extension")
        parser.add_argument("--scan-core", action="store_true", 
                          help="Scan core script directories (projectiles, base_classes, enemies)")
        
        args = parser.parse_args()
        
        # Special mode: Scan core scripts
        if args.scan_core:
            scan_core_scripts(args.directory, args.output)
            return 0
        
        # Set filter mode based on arguments
        filter_mode = FilterMode.INCLUDE if args.include else FilterMode.EXCLUDE
        
        # Create a config with default options
        config = create_default_config(args.directory, args.output, filter_mode)
        
        # Update extensions based on command-line argument if provided
        if args.extensions:
            ext_list = [ext.strip() for ext in args.extensions.split(",")]
            # Ensure all extensions have a dot prefix
            ext_list = [ext if ext.startswith(".") else f".{ext}" for ext in ext_list]
            
            if filter_mode == FilterMode.INCLUDE:
                config.included_extensions = set(ext_list)
            else:  # EXCLUDE mode
                config.excluded_extensions = set(ext_list)
                
        # Update excluded files if provided
        if args.exclude_files:
            excluded_files = [file.strip() for file in args.exclude_files.split(",")]
            config.excluded_files = set(excluded_files)
            
        # Update included files if provided
        if args.include_files:
            included_files = [file.strip() for file in args.include_files.split(",")]
            config.included_files = set(included_files)
        
        # Create and run the scanner
        print(f"Scanning directory: {config.directory.absolute()}")
        print(f"Output file: {config.output_file_path.absolute()}")
        print(f"Filter mode: {config.filter_mode.value}")
        
        if config.filter_mode == FilterMode.INCLUDE:
            print(f"Including extensions: {', '.join(sorted(config.included_extensions))}")
        else:
            print(f"Excluding extensions: {', '.join(sorted(config.excluded_extensions))}")
            
        if config.excluded_files:
            print(f"Excluding specific files: {', '.join(sorted(config.excluded_files))}")
            
        if config.included_files:
            print(f"Including specific files: {', '.join(sorted(config.included_files))}")
            
        print("Scanning in progress...")
        
        scanner = Scanner(config)
        scanner.scan()
        
        print(f"Scan completed. Report saved to {config.output_file_path.absolute()}")
        return 0
        
    except KeyboardInterrupt:
        print("\nScan interrupted by user.")
        return 130
    except ScannerConfigError as e:
        print(f"Configuration error: {str(e)}")
        return 1
    except ScannerError as e:
        print(f"Scanning error: {str(e)}")
        return 1
    except Exception as e:
        print(f"Unexpected error: {str(e)}")
        return 1


if __name__ == "__main__":
    sys.exit(main())